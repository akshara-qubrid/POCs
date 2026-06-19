"""
pptx_builder.py — Modern, visually premium slide generation system.

Inspired by Gamma, Tome, Beautiful.ai, Linear, Notion, and Anthropic design.
Generates 7–9 slides with dedicated renderer per slide type:
  CoverSlide, ProblemSlide, SolutionSlide, ProductSlide,
  ArchitectureSlide, ComparisonSlide, RoadmapSlide, ClosingSlide

Design principles:
  - Large typography, strong visual hierarchy
  - Rounded cards, consistent spacing, generous whitespace
  - Feature cards, KPI/stat cards, two-column layouts
  - Max 5–6 content blocks per slide
  - Premium SaaS aesthetic (dark + vivid accent palette)

Slide dimensions: 16:9 widescreen (13.33" × 7.5")
"""

from __future__ import annotations

import io
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
import lxml.etree as etree

# ---------------------------------------------------------------------------
# DESIGN TOKENS
# ---------------------------------------------------------------------------

# Base palette
C_BG_DEEP   = RGBColor(0x09, 0x0C, 0x14)   # near-black canvas
C_BG_CARD   = RGBColor(0x13, 0x17, 0x24)   # card surface
C_BG_CARD2  = RGBColor(0x1C, 0x21, 0x35)   # elevated card
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_SUBTEXT   = RGBColor(0x8B, 0x95, 0xB4)   # muted label
C_DIVIDER   = RGBColor(0x25, 0x2C, 0x46)   # subtle separator

# Slide-type accent palette  (index matches SLIDE_ORDER)
ACCENTS: List[Tuple[int, int, int]] = [
    (0xF5, 0x9E, 0x0B),   # 0 Cover        – amber
    (0xF8, 0x71, 0x71),   # 1 Problem       – coral-red
    (0x34, 0xD3, 0x99),   # 2 Solution      – emerald
    (0x60, 0xA5, 0xFA),   # 3 Product       – sky-blue
    (0xA7, 0x8B, 0xFA),   # 4 Architecture  – violet
    (0xFB, 0xD3, 0x8D),   # 5 Comparison    – peach-gold
    (0x38, 0xBD, 0xF8),   # 6 Roadmap       – cyan
    (0xF5, 0x9E, 0x0B),   # 7 Closing       – amber (bookend)
]

# Slide dimension constants
_W = Inches(13.333)
_H = Inches(7.5)

# Layout guides
_MARGIN_X = Inches(0.65)
_MARGIN_Y = Inches(0.55)
_CONTENT_W = _W - 2 * _MARGIN_X


# ---------------------------------------------------------------------------
# LOW-LEVEL HELPERS
# ---------------------------------------------------------------------------

def _rgb(r: int, g: int, b: int) -> RGBColor:
    return RGBColor(r, g, b)


def _acc(t: Tuple[int, int, int]) -> RGBColor:
    """Unpack accent tuple to RGBColor."""
    return RGBColor(*t)


def _solid_bg(slide, r: int, g: int, b: int) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(r, g, b)


def _rect(slide, left, top, w, h,
          color: Tuple[int, int, int],
          radius_emu: int = 0) -> Any:
    """Add a filled rectangle with optional rounded corners (via XML)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(*color)
    shape.line.fill.background()
    if radius_emu:
        # Set rounded corners via prstGeom adj
        sp = shape._element
        pPr = sp.find(qn("p:spPr"))
        if pPr is not None:
            prstGeom = pPr.find(qn("a:prstGeom"))
            if prstGeom is not None:
                pPr.remove(prstGeom)
            rnd = etree.SubElement(pPr, qn("a:prstGeom"))
            rnd.set("prst", "roundRect")
            avLst = etree.SubElement(rnd, qn("a:avLst"))
            gd = etree.SubElement(avLst, qn("a:gd"))
            gd.set("name", "adj")
            # adj value: 0–50000, proportional to min(w,h)
            min_dim = min(w, h)
            adj_val = min(50000, int(radius_emu / min_dim * 100000)) if min_dim else 10000
            gd.set("fmla", f"val {adj_val}")
    return shape


def _textbox(slide, text: str,
             left, top, w, h,
             size: int = 18,
             bold: bool = False,
             color: RGBColor = None,
             align: Any = PP_ALIGN.LEFT,
             wrap: bool = True,
             italic: bool = False,
             spacing_before: int = 0) -> None:
    if color is None:
        color = C_WHITE
    txb = slide.shapes.add_textbox(left, top, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing_before:
        p.space_before = Pt(spacing_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"


def _multiline_textbox(slide, lines: List[str],
                       left, top, w, h,
                       size: int = 16,
                       color: RGBColor = None,
                       bullet_color: RGBColor = None,
                       bullet_char: str = "▸",
                       line_spacing: int = 6) -> None:
    """Render a list of strings as bullet lines inside a textbox."""
    if color is None:
        color = C_WHITE
    if bullet_color is None:
        bullet_color = color
    txb = slide.shapes.add_textbox(left, top, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_before = Pt(line_spacing)
        if bullet_char:
            m = p.add_run()
            m.text = f"{bullet_char}  "
            m.font.size = Pt(size)
            m.font.bold = True
            m.font.color.rgb = bullet_color
            m.font.name = "Segoe UI"
        r = p.add_run()
        r.text = str(line)
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Segoe UI"


def _slide_header(slide, title: str, accent: Tuple[int, int, int],
                  slide_num: Optional[int] = None) -> None:
    """Render the top accent bar + slide title + optional slide number."""
    # Thin accent top bar
    _rect(slide, 0, 0, _W, Inches(0.06), accent)

    # Slide number badge
    if slide_num is not None:
        _textbox(slide, f"0{slide_num}" if slide_num < 10 else str(slide_num),
                 _W - Inches(1.0), Inches(0.12), Inches(0.75), Inches(0.35),
                 size=11, color=C_SUBTEXT, align=PP_ALIGN.RIGHT)

    # Title
    _textbox(slide, title,
             _MARGIN_X, Inches(0.22), _CONTENT_W, Inches(0.7),
             size=34, bold=True, color=_acc(accent))

    # Divider line
    _rect(slide, _MARGIN_X, Inches(1.0), _CONTENT_W, Pt(1.5).emu, accent)


def _notes(slide, text: str) -> None:
    if text:
        slide.notes_slide.notes_text_frame.text = str(text)


def _coerce_list(raw: Any) -> List[str]:
    if isinstance(raw, list):
        return [str(x).strip() for x in raw if str(x).strip()]
    if isinstance(raw, str):
        lines = [ln.strip(" -•▸→*") for ln in raw.replace(";", "\n").splitlines()]
        return [ln for ln in lines if ln]
    return []


def _blank_slide(prs: Presentation) -> Any:
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


# ---------------------------------------------------------------------------
# COMPONENT: Feature Card
# ---------------------------------------------------------------------------

def _feature_card(slide,
                  left, top, w, h,
                  title: str, body: str,
                  accent: Tuple[int, int, int],
                  icon: str = "") -> None:
    """Render a rounded card with a coloured accent top-edge, title, and body."""
    CARD_BG = (0x1C, 0x21, 0x35)
    RADIUS = Inches(0.12)

    _rect(slide, left, top, w, h, CARD_BG, radius_emu=int(RADIUS))

    # Accent top-edge strip inside card
    _rect(slide, left, top, w, Inches(0.055), accent)

    # Icon / emoji
    offset_y = Inches(0.18)
    if icon:
        _textbox(slide, icon, left + Inches(0.2), top + offset_y,
                 Inches(0.55), Inches(0.45), size=20, align=PP_ALIGN.LEFT)
        title_left = left + Inches(0.75)
    else:
        title_left = left + Inches(0.2)

    # Card title
    _textbox(slide, title, title_left, top + offset_y,
             w - Inches(0.9), Inches(0.42),
             size=15, bold=True, color=C_WHITE)

    # Card body
    _textbox(slide, body, left + Inches(0.2), top + Inches(0.72),
             w - Inches(0.4), h - Inches(0.82),
             size=12, color=C_SUBTEXT, wrap=True)


# ---------------------------------------------------------------------------
# COMPONENT: KPI / Stat Card
# ---------------------------------------------------------------------------

def _stat_card(slide,
               left, top, w, h,
               metric: str, value: str, subtitle: str,
               accent: Tuple[int, int, int]) -> None:
    """Large number / metric card — highlights a single key stat."""
    CARD_BG = (0x1C, 0x21, 0x35)
    _rect(slide, left, top, w, h, CARD_BG, radius_emu=int(Inches(0.12)))

    # Large accent value
    _textbox(slide, value,
             left + Inches(0.15), top + Inches(0.2), w - Inches(0.3), Inches(0.8),
             size=36, bold=True, color=_acc(accent), align=PP_ALIGN.CENTER)

    # Metric label
    _textbox(slide, metric,
             left + Inches(0.15), top + Inches(1.0), w - Inches(0.3), Inches(0.35),
             size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Sub-label
    if subtitle:
        _textbox(slide, subtitle,
                 left + Inches(0.15), top + Inches(1.35), w - Inches(0.3), Inches(0.3),
                 size=10, color=C_SUBTEXT, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# SLIDE RENDERER: Cover
# ---------------------------------------------------------------------------

def render_cover(prs: Presentation, data: Dict[str, Any],
                 deck_title: str, tagline: str,
                 accent: Tuple[int, int, int]) -> None:
    slide = _blank_slide(prs)
    _solid_bg(slide, *C_BG_DEEP)

    # Full-width gradient overlay strip (decorative)
    _rect(slide, 0, Inches(5.5), _W, Inches(2.0), (0x13, 0x17, 0x24))

    # Left vertical accent bar
    _rect(slide, 0, 0, Inches(0.22), _H, accent)

    # Corner glow — large soft rectangle in top-right
    _rect(slide, _W - Inches(3.5), 0, Inches(3.5), Inches(2.8), (0x18, 0x1A, 0x2A))

    # Accent dot cluster (decorative circles via small rects)
    for i, dy in enumerate([0.3, 0.65, 1.0]):
        a = list(accent)
        dot_color = (
            max(0, min(255, int(a[0] // (i + 1)) + 30)),
            max(0, min(255, int(a[1] // (i + 1)) + 20)),
            max(0, min(255, int(a[2] // (i + 1)) + 10)),
        )
        _rect(slide, _W - Inches(0.55), Inches(dy), Inches(0.12), Inches(0.12), dot_color)

    # Company / deck title — large hero text
    _textbox(slide, deck_title,
             Inches(1.1), Inches(1.8), Inches(10.5), Inches(1.6),
             size=58, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    # Accent underline
    _rect(slide, Inches(1.1), Inches(3.5), Inches(3.5), Inches(0.06), accent)

    # Tagline
    _textbox(slide, tagline,
             Inches(1.1), Inches(3.7), Inches(10.0), Inches(0.7),
             size=22, bold=False, color=_acc(accent), align=PP_ALIGN.LEFT)

    # Meta chips row
    content = _coerce_list(data.get("content", []))
    chips = "   ·   ".join(content[:4])
    if chips:
        _textbox(slide, chips,
                 Inches(1.1), Inches(4.6), Inches(10.0), Inches(0.5),
                 size=13, color=C_SUBTEXT, align=PP_ALIGN.LEFT)

    # Bottom thin bar
    _rect(slide, 0, _H - Inches(0.06), _W, Inches(0.06), accent)

    _notes(slide, data.get("speaker_notes", ""))


# ---------------------------------------------------------------------------
# SLIDE RENDERER: Problem
# ---------------------------------------------------------------------------

def render_problem(prs: Presentation, data: Dict[str, Any],
                   slide_num: int, accent: Tuple[int, int, int]) -> None:
    slide = _blank_slide(prs)
    _solid_bg(slide, *C_BG_DEEP)
    _slide_header(slide, data.get("title", "The Problem"), accent, slide_num)

    items = _coerce_list(data.get("content", []))[:5]
    n = len(items)
    if n == 0:
        _notes(slide, data.get("speaker_notes", ""))
        return

    # Left column: large callout / pain statement
    if items:
        pain = items[0]
        _rect(slide, _MARGIN_X, Inches(1.2), Inches(4.8), Inches(5.6),
              (0x1F, 0x14, 0x14), radius_emu=int(Inches(0.14)))
        # Red-tinted label
        _textbox(slide, "PAIN POINT",
                 _MARGIN_X + Inches(0.25), Inches(1.45), Inches(4.3), Inches(0.35),
                 size=10, bold=True, color=_acc(accent), italic=True)
        _textbox(slide, pain,
                 _MARGIN_X + Inches(0.25), Inches(1.85), Inches(4.3), Inches(4.5),
                 size=20, color=C_WHITE, wrap=True)

    # Right column: supporting problem bullets as cards
    right_x = _MARGIN_X + Inches(5.2)
    right_w = _W - right_x - _MARGIN_X
    remaining = items[1:5]
    card_h = Inches(1.18) if len(remaining) >= 4 else Inches(1.45)
    for i, item in enumerate(remaining):
        cy = Inches(1.2) + i * (card_h + Inches(0.12))
        _rect(slide, right_x, cy, right_w, card_h, (0x1C, 0x21, 0x35),
              radius_emu=int(Inches(0.1)))
        # Small accent circle index
        _textbox(slide, str(i + 1),
                 right_x + Inches(0.2), cy + Inches(0.28), Inches(0.3), Inches(0.4),
                 size=16, bold=True, color=_acc(accent))
        _textbox(slide, item,
                 right_x + Inches(0.6), cy + Inches(0.1), right_w - Inches(0.75), card_h - Inches(0.15),
                 size=15, color=C_WHITE, wrap=True)

    _rect(slide, 0, _H - Inches(0.06), _W, Inches(0.06), accent)
    _notes(slide, data.get("speaker_notes", ""))


# ---------------------------------------------------------------------------
# SLIDE RENDERER: Solution
# ---------------------------------------------------------------------------

def render_solution(prs: Presentation, data: Dict[str, Any],
                    slide_num: int, accent: Tuple[int, int, int]) -> None:
    slide = _blank_slide(prs)
    _solid_bg(slide, *C_BG_DEEP)
    _slide_header(slide, data.get("title", "Our Solution"), accent, slide_num)

    items = _coerce_list(data.get("content", []))[:6]

    ICONS = ["⚡", "🎯", "🔒", "🚀", "💡", "🔗"]
    n = len(items)
    if n == 0:
        _notes(slide, data.get("speaker_notes", ""))
        return

    cols = 3 if n >= 4 else (2 if n >= 2 else 1)
    rows = (n + cols - 1) // cols
    gutter = Inches(0.18)
    card_w = (_CONTENT_W - gutter * (cols - 1)) / cols
    card_h = (Inches(5.6) - gutter * (rows - 1)) / rows
    start_y = Inches(1.2)

    for i, item in enumerate(items):
        row = i // cols
        col = i % cols
        cx = _MARGIN_X + col * (card_w + gutter)
        cy = start_y + row * (card_h + gutter)
        # Split item into title / body if separator present
        if " — " in item:
            parts = item.split(" — ", 1)
            ctitle, cbody = parts[0], parts[1]
        elif ": " in item:
            parts = item.split(": ", 1)
            ctitle, cbody = parts[0], parts[1]
        else:
            ctitle = item[:40]
            cbody = item[40:] if len(item) > 40 else ""
        _feature_card(slide, cx, cy, card_w, card_h,
                      ctitle, cbody, accent, icon=ICONS[i % len(ICONS)])

    _rect(slide, 0, _H - Inches(0.06), _W, Inches(0.06), accent)
    _notes(slide, data.get("speaker_notes", ""))


# ---------------------------------------------------------------------------
# SLIDE RENDERER: Product Overview
# ---------------------------------------------------------------------------

def render_product(prs: Presentation, data: Dict[str, Any],
                   slide_num: int, accent: Tuple[int, int, int]) -> None:
    slide = _blank_slide(prs)
    _solid_bg(slide, *C_BG_DEEP)
    _slide_header(slide, data.get("title", "Product Overview"), accent, slide_num)

    items = _coerce_list(data.get("content", []))[:6]
    stats = data.get("stats", [])  # optional list of {metric, value, subtitle}

    # If stats provided, show KPI row first, then feature list
    stat_list = stats if isinstance(stats, list) else []
    stat_count = min(len(stat_list), 4)
    content_start_y = Inches(1.2)

    if stat_count > 0:
        stat_w = (_CONTENT_W - Inches(0.15) * (stat_count - 1)) / stat_count
        for i, s in enumerate(stat_list[:stat_count]):
            sx = _MARGIN_X + i * (stat_w + Inches(0.15))
            _stat_card(slide, sx, Inches(1.2), stat_w, Inches(1.85),
                       s.get("metric", ""), s.get("value", ""), s.get("subtitle", ""), accent)
        content_start_y = Inches(3.25)

    # Feature list: two columns
    if items:
        half = (len(items) + 1) // 2
        left_items = items[:half]
        right_items = items[half:]
        col_w = (_CONTENT_W - Inches(0.25)) / 2
        avail_h = _H - content_start_y - Inches(0.55)

        for col_idx, col_items in enumerate([left_items, right_items]):
            cx = _MARGIN_X + col_idx * (col_w + Inches(0.25))
            _multiline_textbox(slide, col_items, cx, content_start_y, col_w, avail_h,
                               size=16, color=C_WHITE,
                               bullet_color=_acc(accent), bullet_char="▸")

    _rect(slide, 0, _H - Inches(0.06), _W, Inches(0.06), accent)
    _notes(slide, data.get("speaker_notes", ""))


# ---------------------------------------------------------------------------
# SLIDE RENDERER: Architecture / How It Works
# ---------------------------------------------------------------------------

def render_architecture(prs: Presentation, data: Dict[str, Any],
                         slide_num: int, accent: Tuple[int, int, int]) -> None:
    slide = _blank_slide(prs)
    _solid_bg(slide, *C_BG_DEEP)
    _slide_header(slide, data.get("title", "Architecture"), accent, slide_num)

    items = _coerce_list(data.get("content", []))[:6]
    n = len(items)

    # Render as a horizontal pipeline / step diagram
    STEP_ICONS = ["①", "②", "③", "④", "⑤", "⑥"]
    step_count = min(n, 5)
    if step_count == 0:
        _notes(slide, data.get("speaker_notes", ""))
        return

    box_w = (_CONTENT_W - Inches(0.15) * (step_count - 1)) / step_count
    box_h = Inches(2.4)
    arrow_color = (0x35, 0x3C, 0x58)
    start_y = Inches(2.1)

    for i in range(step_count):
        bx = _MARGIN_X + i * (box_w + Inches(0.15))

        # Step box
        _rect(slide, bx, start_y, box_w, box_h, (0x1C, 0x21, 0x35),
              radius_emu=int(Inches(0.12)))

        # Accent top strip
        _rect(slide, bx, start_y, box_w, Inches(0.06), accent)

        # Step number
        _textbox(slide, STEP_ICONS[i],
                 bx + Inches(0.15), start_y + Inches(0.1), Inches(0.5), Inches(0.5),
                 size=18, bold=True, color=_acc(accent))

        # Step label (title part)
        label = items[i]
        if ": " in label:
            stitle, sbody = label.split(": ", 1)
        elif " — " in label:
            stitle, sbody = label.split(" — ", 1)
        else:
            stitle = label[:30]
            sbody = label[30:]
        _textbox(slide, stitle,
                 bx + Inches(0.15), start_y + Inches(0.65), box_w - Inches(0.25), Inches(0.45),
                 size=13, bold=True, color=C_WHITE, wrap=True)
        if sbody:
            _textbox(slide, sbody,
                     bx + Inches(0.15), start_y + Inches(1.15), box_w - Inches(0.25), Inches(1.1),
                     size=11, color=C_SUBTEXT, wrap=True)

        # Arrow connector between boxes (except last)
        if i < step_count - 1:
            arrow_x = bx + box_w + Inches(0.02)
            _textbox(slide, "→",
                     arrow_x, start_y + box_h / 2 - Inches(0.2), Inches(0.11), Inches(0.4),
                     size=14, color=_acc(accent), align=PP_ALIGN.CENTER)

    # Below pipeline: remaining items as two-col bullets
    if n > step_count:
        extra = items[step_count:]
        _multiline_textbox(slide, extra,
                           _MARGIN_X, start_y + box_h + Inches(0.3),
                           _CONTENT_W, Inches(1.4),
                           size=14, color=C_SUBTEXT, bullet_color=_acc(accent))

    _rect(slide, 0, _H - Inches(0.06), _W, Inches(0.06), accent)
    _notes(slide, data.get("speaker_notes", ""))


# ---------------------------------------------------------------------------
# SLIDE RENDERER: Comparison / Competitive Advantage
# ---------------------------------------------------------------------------

def render_comparison(prs: Presentation, data: Dict[str, Any],
                       slide_num: int, accent: Tuple[int, int, int]) -> None:
    slide = _blank_slide(prs)
    _solid_bg(slide, *C_BG_DEEP)
    _slide_header(slide, data.get("title", "Competitive Advantages"), accent, slide_num)

    items = _coerce_list(data.get("content", []))[:6]

    # Two-column feature advantage layout
    # Left col: "Others" (red — lack) | Right col: "Us" (accent — have)
    col_w = (_CONTENT_W - Inches(0.3)) / 2
    left_x  = _MARGIN_X
    right_x = _MARGIN_X + col_w + Inches(0.3)
    start_y = Inches(1.25)

    # Column headers
    _rect(slide, left_x, start_y, col_w, Inches(0.5), (0x2A, 0x14, 0x14),
          radius_emu=int(Inches(0.08)))
    _textbox(slide, "Traditional Alternatives",
             left_x + Inches(0.2), start_y + Inches(0.08), col_w - Inches(0.3), Inches(0.36),
             size=13, bold=True, color=_rgb(0xF8, 0x71, 0x71), align=PP_ALIGN.CENTER)

    _rect(slide, right_x, start_y, col_w, Inches(0.5), (0x0F, 0x28, 0x20),
          radius_emu=int(Inches(0.08)))
    _textbox(slide, "Our Approach",
             right_x + Inches(0.2), start_y + Inches(0.08), col_w - Inches(0.3), Inches(0.36),
             size=13, bold=True, color=_acc(accent), align=PP_ALIGN.CENTER)

    # Comparison rows
    row_h = Inches(0.7)
    row_start = start_y + Inches(0.65)
    for i, item in enumerate(items):
        ry = row_start + i * (row_h + Inches(0.1))
        # Alternating row background
        row_bg = (0x18, 0x1C, 0x2E) if i % 2 == 0 else (0x1C, 0x21, 0x35)

        # Left cell — problem / lack
        _rect(slide, left_x, ry, col_w, row_h, row_bg, radius_emu=int(Inches(0.07)))
        _textbox(slide, "✗  " + item,
                 left_x + Inches(0.15), ry + Inches(0.1), col_w - Inches(0.25), row_h - Inches(0.15),
                 size=13, color=_rgb(0xF8, 0x71, 0x71), wrap=True)

        # Right cell — our advantage
        advantage = data.get("advantages", {}).get(str(i), f"We solve: {item[:60]}")
        _rect(slide, right_x, ry, col_w, row_h, row_bg, radius_emu=int(Inches(0.07)))
        _textbox(slide, "✓  " + advantage,
                 right_x + Inches(0.15), ry + Inches(0.1), col_w - Inches(0.25), row_h - Inches(0.15),
                 size=13, color=_acc(accent), wrap=True)

    _rect(slide, 0, _H - Inches(0.06), _W, Inches(0.06), accent)
    _notes(slide, data.get("speaker_notes", ""))


# ---------------------------------------------------------------------------
# SLIDE RENDERER: Use Cases / Market Opportunity
# ---------------------------------------------------------------------------

def render_usecases(prs: Presentation, data: Dict[str, Any],
                    slide_num: int, accent: Tuple[int, int, int]) -> None:
    slide = _blank_slide(prs)
    _solid_bg(slide, *C_BG_DEEP)
    _slide_header(slide, data.get("title", "Use Cases & Market Opportunity"), accent, slide_num)

    items = _coerce_list(data.get("content", []))[:6]
    stats = data.get("stats", [])
    stat_list = stats if isinstance(stats, list) else []

    # Top row: up to 3 KPI stat cards for TAM/SAM/SOM or key market metrics
    if stat_list:
        sc = min(len(stat_list), 3)
        sw = (_CONTENT_W - Inches(0.2) * (sc - 1)) / sc
        for i, s in enumerate(stat_list[:sc]):
            sx = _MARGIN_X + i * (sw + Inches(0.2))
            _stat_card(slide, sx, Inches(1.2), sw, Inches(1.7),
                       s.get("metric", ""), s.get("value", ""), s.get("subtitle", ""), accent)
        card_start_y = Inches(3.1)
    else:
        card_start_y = Inches(1.2)

    # Bottom: use case cards (2 or 3 columns)
    avail_h = _H - card_start_y - Inches(0.55)
    n = len(items)
    if n == 0:
        _rect(slide, 0, _H - Inches(0.06), _W, Inches(0.06), accent)
        _notes(slide, data.get("speaker_notes", ""))
        return

    cols = 3 if n >= 4 else (2 if n >= 2 else 1)
    rows = (n + cols - 1) // cols
    gutter = Inches(0.15)
    card_w = (_CONTENT_W - gutter * (cols - 1)) / cols
    card_h = (avail_h - gutter * (rows - 1)) / rows
    USE_ICONS = ["🏢", "👥", "⚙️", "📊", "🔄", "🌐"]

    for i, item in enumerate(items):
        row = i // cols
        col = i % cols
        cx = _MARGIN_X + col * (card_w + gutter)
        cy = card_start_y + row * (card_h + gutter)
        if ": " in item:
            t, b = item.split(": ", 1)
        elif " — " in item:
            t, b = item.split(" — ", 1)
        else:
            t = item[:35]
            b = item[35:]
        _feature_card(slide, cx, cy, card_w, card_h, t, b, accent, USE_ICONS[i % len(USE_ICONS)])

    _rect(slide, 0, _H - Inches(0.06), _W, Inches(0.06), accent)
    _notes(slide, data.get("speaker_notes", ""))


# ---------------------------------------------------------------------------
# SLIDE RENDERER: Roadmap
# ---------------------------------------------------------------------------

def render_roadmap(prs: Presentation, data: Dict[str, Any],
                   slide_num: int, accent: Tuple[int, int, int]) -> None:
    slide = _blank_slide(prs)
    _solid_bg(slide, *C_BG_DEEP)
    _slide_header(slide, data.get("title", "Roadmap"), accent, slide_num)

    items = _coerce_list(data.get("content", []))[:6]
    n = len(items)
    if n == 0:
        _rect(slide, 0, _H - Inches(0.06), _W, Inches(0.06), accent)
        _notes(slide, data.get("speaker_notes", ""))
        return

    # Horizontal timeline track
    track_y = Inches(2.5)
    track_h = Inches(0.05)
    _rect(slide, _MARGIN_X, track_y + Inches(0.72), _CONTENT_W, track_h,
          (0x35, 0x3C, 0x58))

    # Phase items spaced evenly along track
    spacing = _CONTENT_W / n
    PHASE_LABELS = ["Now", "Q2", "Q3", "Q4", "2026", "2027"]
    BOX_W = min(Inches(2.3), spacing - Inches(0.15))
    BOX_H = Inches(1.5)

    for i, item in enumerate(items):
        cx = _MARGIN_X + i * spacing + (spacing - BOX_W) / 2
        cy = Inches(1.2) if i % 2 == 0 else Inches(3.8)
        dot_y = track_y + Inches(0.72) - Inches(0.015)
        dot_x = _MARGIN_X + i * spacing + spacing / 2 - Inches(0.12)

        # Connector line from box to dot
        line_x = dot_x + Inches(0.12)
        if i % 2 == 0:
            # Box above → line goes down
            _rect(slide, line_x, cy + BOX_H, Pt(1.5).emu, dot_y - cy - BOX_H,
                  (0x35, 0x3C, 0x58))
        else:
            # Box below → line goes up
            _rect(slide, line_x, dot_y + Inches(0.025), Pt(1.5).emu, cy - dot_y,
                  (0x35, 0x3C, 0x58))

        # Dot on track
        _rect(slide, dot_x, dot_y, Inches(0.22), Inches(0.22), accent,
              radius_emu=int(Inches(0.11)))

        # Phase label below dot
        _textbox(slide, PHASE_LABELS[i] if i < len(PHASE_LABELS) else f"P{i+1}",
                 dot_x - Inches(0.3), dot_y + Inches(0.28), Inches(0.8), Inches(0.3),
                 size=10, bold=True, color=_acc(accent), align=PP_ALIGN.CENTER)

        # Card box
        _rect(slide, cx, cy, BOX_W, BOX_H, (0x1C, 0x21, 0x35),
              radius_emu=int(Inches(0.1)))
        _rect(slide, cx, cy, BOX_W, Inches(0.05), accent)
        _textbox(slide, item, cx + Inches(0.15), cy + Inches(0.12),
                 BOX_W - Inches(0.25), BOX_H - Inches(0.2),
                 size=12, color=C_WHITE, wrap=True)

    _rect(slide, 0, _H - Inches(0.06), _W, Inches(0.06), accent)
    _notes(slide, data.get("speaker_notes", ""))


# ---------------------------------------------------------------------------
# SLIDE RENDERER: Closing / CTA
# ---------------------------------------------------------------------------

def render_closing(prs: Presentation, data: Dict[str, Any],
                   deck_title: str, accent: Tuple[int, int, int]) -> None:
    slide = _blank_slide(prs)
    _solid_bg(slide, *C_BG_DEEP)

    # Bottom accent bar
    _rect(slide, 0, _H - Inches(0.06), _W, Inches(0.06), accent)
    _rect(slide, 0, 0, Inches(0.22), _H, accent)

    # Large CTA text
    _textbox(slide, data.get("title", "Let's Build the Future Together"),
             Inches(1.1), Inches(1.6), Inches(11.0), Inches(1.5),
             size=46, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    # Accent underline
    _rect(slide, Inches(1.1), Inches(3.25), Inches(4.5), Inches(0.07), accent)

    # Subtitle / contact prompt
    items = _coerce_list(data.get("content", []))
    subtitle = items[0] if items else "We'd love to explore this with you."
    _textbox(slide, subtitle,
             Inches(1.1), Inches(3.45), Inches(10.0), Inches(0.7),
             size=20, color=_acc(accent))

    # Contact / CTA line items
    for i, item in enumerate(items[1:4]):
        _textbox(slide, f"  {item}",
                 Inches(1.1), Inches(4.35) + i * Inches(0.5), Inches(9.0), Inches(0.45),
                 size=15, color=C_SUBTEXT)

    # Company name watermark bottom-right
    _textbox(slide, deck_title,
             _W - Inches(4.5), _H - Inches(0.55), Inches(4.0), Inches(0.4),
             size=12, color=C_SUBTEXT, align=PP_ALIGN.RIGHT)

    _notes(slide, data.get("speaker_notes", ""))


# ---------------------------------------------------------------------------
# SLIDE TYPE REGISTRY
# ---------------------------------------------------------------------------

# Maps the 'slide_type' field (or title keywords) to a renderer function name
_SLIDE_TYPE_MAP: Dict[str, str] = {
    "cover":          "cover",
    "problem":        "problem",
    "solution":       "solution",
    "product":        "product",
    "architecture":   "architecture",
    "how it works":   "architecture",
    "use cases":      "usecases",
    "market":         "usecases",
    "competitive":    "comparison",
    "competition":    "comparison",
    "advantages":     "comparison",
    "roadmap":        "roadmap",
    "closing":        "closing",
    "conclusion":     "closing",
    "contact":        "closing",
}

def _detect_slide_type(slide_data: Dict[str, Any], idx: int) -> str:
    """Determine renderer type from explicit 'slide_type' field or title heuristics."""
    explicit = slide_data.get("slide_type", "").lower().strip()
    if explicit in _SLIDE_TYPE_MAP:
        return _SLIDE_TYPE_MAP[explicit]
    # Fallback: match by title keywords
    title = slide_data.get("title", "").lower()
    for keyword, rtype in _SLIDE_TYPE_MAP.items():
        if keyword in title:
            return rtype
    # Index-based fallback for typical deck order
    ORDER = ["cover", "problem", "solution", "product",
             "architecture", "usecases", "comparison", "roadmap", "closing"]
    return ORDER[min(idx, len(ORDER) - 1)]


# ---------------------------------------------------------------------------
# PUBLIC API: build_pptx
# ---------------------------------------------------------------------------

def build_pptx(deck: Dict[str, Any]) -> bytes:
    """
    Convert a structured deck dict into a premium .pptx file.
    Returns raw bytes for streaming as a file download.

    Expected deck structure:
        {
          "title": "Startup Name",
          "tagline": "One-line value prop",
          "slides": [
            {
              "slide_number": 1,
              "slide_type": "cover",       # optional — auto-detected if absent
              "title": "...",
              "content": ["...", ...],     # list of strings, max 6
              "stats": [...],              # optional KPI cards
              "advantages": {...},         # optional for comparison slide
              "speaker_notes": "..."
            }, ...
          ]
        }
    """
    prs = Presentation()
    prs.slide_width  = _W
    prs.slide_height = _H

    deck_title = deck.get("title", "Pitch Deck")
    tagline    = deck.get("tagline", "")
    slides: List[Dict[str, Any]] = deck.get("slides", [])

    for i, slide_data in enumerate(slides):
        acc = ACCENTS[i % len(ACCENTS)]
        stype = _detect_slide_type(slide_data, i)

        if stype == "cover":
            render_cover(prs, slide_data, deck_title, tagline, acc)
        elif stype == "problem":
            render_problem(prs, slide_data, i + 1, acc)
        elif stype == "solution":
            render_solution(prs, slide_data, i + 1, acc)
        elif stype == "product":
            render_product(prs, slide_data, i + 1, acc)
        elif stype == "architecture":
            render_architecture(prs, slide_data, i + 1, acc)
        elif stype == "usecases":
            render_usecases(prs, slide_data, i + 1, acc)
        elif stype == "comparison":
            render_comparison(prs, slide_data, i + 1, acc)
        elif stype == "roadmap":
            render_roadmap(prs, slide_data, i + 1, acc)
        elif stype == "closing":
            render_closing(prs, slide_data, deck_title, acc)
        else:
            # Generic fallback — render as product-style slide
            render_product(prs, slide_data, i + 1, acc)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
