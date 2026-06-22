"""
render_pptx.py — Stage 4 of the 4-stage pitch deck pipeline.

Pure python-pptx shape factory. Zero LLM involvement. Zero design decisions.

Responsibilities:
  - Accept a list of RenderedSlide objects (from Stage 3) and a Theme.
  - For each slide, set the background and iterate shapes in z_order.
  - Translate each ShapeSpec into the corresponding python-pptx API call.
  - All font sizes, colors, and geometry come directly from the ShapeSpec —
    this stage never makes any of those choices itself.

If something looks wrong visually, the bug lives in layout_engine.py, not here.
"""
from __future__ import annotations

import io
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.chart.data import ChartData as PptxChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
import lxml.etree as etree

from .themes import Theme
from .layout_engine import RenderedSlide, ShapeSpec

# ---------------------------------------------------------------------------
# Dimension constants (must match layout_engine.py)
# ---------------------------------------------------------------------------

_W = Inches(13.333)
_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------


def _rgb(hex_str: str) -> RGBColor:
    """Convert a 6-char hex string to RGBColor."""
    h = hex_str.strip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _inches(v: float) -> Emu:
    return Inches(v)


def _pt(v: Optional[int]) -> Optional[Emu]:
    return Pt(v) if v is not None else None


def _align(align_str: str) -> Any:
    mapping = {
        "left":   PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right":  PP_ALIGN.RIGHT,
    }
    return mapping.get(align_str, PP_ALIGN.LEFT)


def _apply_rounded_corners(shape: Any, corner_radius_in: float = 0.1) -> None:
    """Apply rounded corners to a shape via XML adjustment values."""
    try:
        sp = shape._element
        pPr = sp.find(qn("p:spPr"))
        if pPr is None:
            return
        prstGeom = pPr.find(qn("a:prstGeom"))
        if prstGeom is not None:
            pPr.remove(prstGeom)
        rnd = etree.SubElement(pPr, qn("a:prstGeom"))
        rnd.set("prst", "roundRect")
        avLst = etree.SubElement(rnd, qn("a:avLst"))
        gd = etree.SubElement(avLst, qn("a:gd"))
        gd.set("name", "adj")
        w = shape.width or 1
        h = shape.height or 1
        min_dim = min(w, h)
        radius_emu = Inches(corner_radius_in)
        adj_val = min(50000, int(radius_emu / min_dim * 100000)) if min_dim else 8000
        gd.set("fmla", f"val {adj_val}")
    except Exception:
        pass  # Silently skip if XML manipulation fails


def _blank_slide(prs: Presentation) -> Any:
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def _set_background(slide: Any, hex_color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex_color)


# ---------------------------------------------------------------------------
# Shape renderers — one per shape_kind
# ---------------------------------------------------------------------------


def _render_rounded_rect(slide: Any, spec: ShapeSpec, theme: Theme) -> None:
    """Render a rounded_rect shape. Used for cards, motif elements, backgrounds."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa — import for side effects
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        _inches(spec.x_in), _inches(spec.y_in),
        _inches(spec.w_in), _inches(spec.h_in),
    )
    if spec.fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(spec.fill_color)
    else:
        shape.fill.background()
    shape.line.fill.background()

    _apply_rounded_corners(shape, theme.motif.corner_radius_in)

    # If this is a motif element, ensure it has no text frame text
    # (motif shapes are purely decorative — their content dict is metadata)


def _render_textbox(slide: Any, spec: ShapeSpec, theme: Theme) -> None:
    """Render a textbox with a single paragraph."""
    c = spec.content
    text = str(c.get("text", ""))
    bold = bool(c.get("bold", False))
    italic = bool(c.get("italic", False))
    align_str = c.get("align", "left")

    txb = slide.shapes.add_textbox(
        _inches(spec.x_in), _inches(spec.y_in),
        _inches(spec.w_in), _inches(spec.h_in),
    )
    tf = txb.text_frame
    tf.word_wrap = True

    # Multi-line text: split on \n
    lines = text.split("\n")
    for line_idx, line in enumerate(lines):
        p = tf.paragraphs[0] if line_idx == 0 else tf.add_paragraph()
        p.alignment = _align(align_str)
        run = p.add_run()
        run.text = line
        if spec.font_size_pt:
            run.font.size = Pt(spec.font_size_pt)
        run.font.bold = bold
        run.font.italic = italic
        if spec.text_color:
            run.font.color.rgb = _rgb(spec.text_color)
        # Apply theme font
        run.font.name = theme.fonts.heading if bold else theme.fonts.body


def _render_table(slide: Any, spec: ShapeSpec, theme: Theme) -> None:
    """Render a python-pptx table from spec.content['table']."""
    table_data = spec.content.get("table", {})
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])

    if not headers and not rows:
        return

    n_cols = max(len(headers), max((len(r) for r in rows), default=0))
    n_rows = len(rows) + (1 if headers else 0)
    if n_cols == 0 or n_rows == 0:
        return

    table = slide.shapes.add_table(
        n_rows, n_cols,
        _inches(spec.x_in), _inches(spec.y_in),
        _inches(spec.w_in), _inches(spec.h_in),
    ).table

    # Header row styling
    if headers:
        for col_i, header_text in enumerate(headers[:n_cols]):
            cell = table.cell(0, col_i)
            cell.text = str(header_text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(theme.colors.primary)
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.bold = True
                    run.font.color.rgb = _rgb(theme.colors.text_on_dark)
                    run.font.name = theme.fonts.heading
                    if spec.font_size_pt:
                        run.font.size = Pt(spec.font_size_pt)

    # Data rows
    row_offset = 1 if headers else 0
    for row_i, row in enumerate(rows):
        for col_i, cell_text in enumerate(row[:n_cols]):
            cell = table.cell(row_i + row_offset, col_i)
            cell.text = str(cell_text)
            # Alternating row background
            fill_hex = theme.colors.card_bg if row_i % 2 == 0 else "FFFFFF"
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(fill_hex)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = theme.fonts.body
                    if spec.font_size_pt:
                        run.font.size = Pt(spec.font_size_pt)
                    if spec.text_color:
                        run.font.color.rgb = _rgb(spec.text_color)


_CHART_TYPE_MAP = {
    "bar":  XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie":  XL_CHART_TYPE.PIE,
}


def _render_chart(slide: Any, spec: ShapeSpec, theme: Theme) -> None:
    """Render a python-pptx chart from spec.content['chart_data']."""
    chart_info = spec.content.get("chart_data", {})
    chart_type_str = chart_info.get("chart_type", "bar")
    chart_type = _CHART_TYPE_MAP.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)
    categories = chart_info.get("categories", [])
    series_list = chart_info.get("series", [])

    if not categories or not series_list:
        # Fallback to image placeholder
        _render_image_placeholder(slide, spec, theme)
        return

    chart_data = PptxChartData()
    chart_data.categories = categories
    for s in series_list:
        chart_data.add_series(s.get("name", ""), s.get("values", []))

    chart = slide.shapes.add_chart(
        chart_type,
        _inches(spec.x_in), _inches(spec.y_in),
        _inches(spec.w_in), _inches(spec.h_in),
        chart_data,
    ).chart

    # Apply theme accent color to the first series
    try:
        if chart.series:
            for i, series in enumerate(chart.series):
                colors = [theme.colors.primary, theme.colors.secondary, theme.colors.accent]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = _rgb(colors[i % len(colors)])
    except Exception:
        pass  # Chart styling is best-effort

    # Style the plot area and chart area backgrounds
    try:
        chart.chart_area.fill.background()
        chart.plot_area.fill.background()
    except Exception:
        pass


def _render_image_placeholder(slide: Any, spec: ShapeSpec, theme: Theme) -> None:
    """
    Render an image placeholder box — a rounded rect with a label.
    Used for team photos and when no image file is available.
    """
    shape = slide.shapes.add_shape(
        1,
        _inches(spec.x_in), _inches(spec.y_in),
        _inches(spec.w_in), _inches(spec.h_in),
    )
    fill_hex = spec.fill_color or theme.colors.primary
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.fill.background()
    _apply_rounded_corners(shape, theme.motif.corner_radius_in)

    label = spec.content.get("label", "")
    if label:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(11)
        run.font.color.rgb = _rgb(theme.colors.text_on_dark)
        run.font.name = theme.fonts.body


# ---------------------------------------------------------------------------
# Shape dispatcher
# ---------------------------------------------------------------------------


_SHAPE_RENDERERS = {
    "rounded_rect":      _render_rounded_rect,
    "textbox":           _render_textbox,
    "table":             _render_table,
    "chart":             _render_chart,
    "image_placeholder": _render_image_placeholder,
}


def _render_shape(slide: Any, spec: ShapeSpec, theme: Theme) -> None:
    renderer = _SHAPE_RENDERERS.get(spec.shape_kind)
    if renderer is None:
        print(f"  [Renderer] WARNING: Unknown shape_kind '{spec.shape_kind}' — skipped.")
        return
    try:
        renderer(slide, spec, theme)
    except Exception as exc:
        print(f"  [Renderer] ERROR rendering {spec.shape_kind} at ({spec.x_in:.2f}, {spec.y_in:.2f}): {exc}")


# ---------------------------------------------------------------------------
# Stage 4 entry point
# ---------------------------------------------------------------------------


def render_deck(
    slides: List[RenderedSlide],
    theme: Theme,
    output_path: Optional[str] = None,
) -> bytes:
    """
    Stage 4 entry point.

    Converts a list of RenderedSlide objects into a .pptx file.

    Args:
        slides:      List of RenderedSlide objects from Stage 3.
        theme:       Theme object from Stage 2 (for font/color application).
        output_path: If provided, also writes the file to disk.

    Returns:
        The .pptx file content as bytes.
    """
    print(f"[Renderer] Stage 4 — rendering {len(slides)} slides to PPTX")

    prs = Presentation()
    prs.slide_width = _W
    prs.slide_height = _H

    for slide_idx, rendered_slide in enumerate(slides):
        pptx_slide = _blank_slide(prs)
        _set_background(pptx_slide, rendered_slide.background_color)

        # Render shapes in z_order (ascending = back to front)
        sorted_shapes = sorted(rendered_slide.shapes, key=lambda s: s.z_order)
        for spec in sorted_shapes:
            _render_shape(pptx_slide, spec, theme)

        print(f"  [Renderer] Slide {slide_idx+1:02d}: {rendered_slide.slide_type} "
              f"({len(rendered_slide.shapes)} shapes)")

    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    if output_path:
        with open(output_path, "wb") as f:
            f.write(pptx_bytes)
        print(f"  [Renderer] Saved to {output_path} ({len(pptx_bytes):,} bytes)")

    return pptx_bytes
